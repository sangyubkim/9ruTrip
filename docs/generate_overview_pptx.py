#!/usr/bin/env python3
"""Generate 9ruTrip product overview PowerPoint (widescreen 13.333x7.5).

Usage:
  python docs/generate_overview_pptx.py
  C:\\Users\\Intellian\\AppData\\Local\\Programs\\Python\\Python312\\python.exe docs\\generate_overview_pptx.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

# --- Brand (from apps/mobile screens) ---
C_PRIMARY = RGBColor(0x0C, 0x4A, 0x6E)  # #0c4a6e
C_ACCENT = RGBColor(0x03, 0x69, 0xA1)  # #0369a1
C_INK = RGBColor(0x0F, 0x17, 0x2A)  # #0f172a
C_BG = RGBColor(0xF8, 0xFA, 0xFC)  # #f8fafc
C_MUTED = RGBColor(0x64, 0x74, 0x8B)  # #64748b
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD = RGBColor(0xFF, 0xFF, 0xFF)
C_LINE = RGBColor(0xE2, 0xE8, 0xF0)
C_SOFT = RGBColor(0xE0, 0xF2, 0xFE)
C_OK = RGBColor(0x04, 0x78, 0x57)

FONT = "Malgun Gothic"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT = Path(__file__).resolve().parent / "9ruTrip-제품개요.pptx"


def _set_run(run, text: str, size: int, bold: bool = False, color: RGBColor = C_INK):
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # Force East Asian font for Malgun Gothic on Windows
    rPr = run._r.get_or_add_rPr()
    ea = parse_xml(
        f'<a:ea xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="{FONT}"/>'
    )
    # remove existing ea if any
    for child in list(rPr):
        if child.tag == "{http://schemas.openxmlformats.org/drawingml/2006/main}ea":
            rPr.remove(child)
    rPr.append(ea)


def add_textbox(slide, left, top, width, height, text, size=14, bold=False, color=C_INK, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        anchor = {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[valign]
        tf._txBody.bodyPr.set("anchor", anchor)
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.clear()
    run = p.add_run()
    _set_run(run, text, size, bold, color)
    return box


def set_para(paragraph, text, size=13, bold=False, color=C_INK, align=PP_ALIGN.LEFT, space_after=6):
    paragraph.alignment = align
    paragraph.space_after = Pt(space_after)
    paragraph.clear()
    run = paragraph.add_run()
    _set_run(run, text, size, bold, color)


def add_bg(slide, color: RGBColor = C_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # send to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return shape


def add_header_bar(slide, title: str, subtitle: str | None = None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_PRIMARY
    bar.line.fill.background()
    add_textbox(slide, Inches(0.45), Inches(0.18), Inches(10), Inches(0.4), title, 22, True, C_WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.45), Inches(0.55), Inches(11), Inches(0.3), subtitle, 11, False, RGBColor(0xBA, 0xE6, 0xFD))
    # brand chip
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.2), Inches(0.28), Inches(1.7), Inches(0.4))
    chip.fill.solid()
    chip.fill.fore_color.rgb = C_ACCENT
    chip.line.fill.background()
    add_textbox(slide, Inches(11.2), Inches(0.32), Inches(1.7), Inches(0.35), "9ruTrip", 12, True, C_WHITE, PP_ALIGN.CENTER)


def add_card(slide, left, top, width, height, fill=C_CARD, line=C_LINE):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = line
    card.line.width = Pt(1)
    return card


def add_bullets(slide, left, top, width, height, items: list[str], size=13, color=C_INK, bold_first=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(5)
        set_para(p, f"•  {item}", size, bold=(bold_first and i == 0), color=color, space_after=5)
    return box


def wireframe_phone(slide, left, top, title: str, lines: list[str], accent_header=True):
    """Simple phone-frame wireframe card."""
    w, h = Inches(3.7), Inches(5.35)
    frame = add_card(slide, left, top, w, h, C_WHITE, C_LINE)
    # status / header
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Inches(0.55))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = C_PRIMARY if accent_header else C_ACCENT
    hdr.line.fill.background()
    add_textbox(slide, left + Inches(0.12), top + Inches(0.12), w - Inches(0.2), Inches(0.35), title, 12, True, C_WHITE)
    y = top + Inches(0.7)
    for line in lines:
        is_btn = line.startswith("[") and line.endswith("]")
        is_sec = line.startswith("##")
        text = line[2:].strip() if is_sec else (line[1:-1] if is_btn else line)
        if is_btn:
            btn = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left + Inches(0.25),
                y,
                w - Inches(0.5),
                Inches(0.38),
            )
            btn.fill.solid()
            btn.fill.fore_color.rgb = C_ACCENT
            btn.line.fill.background()
            add_textbox(
                slide,
                left + Inches(0.25),
                y + Inches(0.05),
                w - Inches(0.5),
                Inches(0.3),
                text,
                10,
                True,
                C_WHITE,
                PP_ALIGN.CENTER,
            )
            y += Inches(0.48)
        elif is_sec:
            add_textbox(slide, left + Inches(0.2), y, w - Inches(0.35), Inches(0.28), text, 10, True, C_PRIMARY)
            y += Inches(0.3)
        elif line.startswith("---"):
            rule = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left + Inches(0.2),
                y + Inches(0.05),
                w - Inches(0.4),
                Pt(1),
            )
            rule.fill.solid()
            rule.fill.fore_color.rgb = C_LINE
            rule.line.fill.background()
            y += Inches(0.2)
        else:
            add_textbox(slide, left + Inches(0.2), y, w - Inches(0.35), Inches(0.28), text, 10, False, C_MUTED)
            y += Inches(0.28)
    return frame


def new_slide(prs: Presentation):
    blank = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(blank)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ========== 1. Cover ==========
    s = new_slide(prs)
    add_bg(s, C_PRIMARY)
    # accent stripe
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.9), SLIDE_W, Inches(1.6))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = RGBColor(0x07, 0x58, 0x85)
    stripe.line.fill.background()
    add_textbox(s, Inches(0.7), Inches(1.8), Inches(11), Inches(0.6), "9ruTrip", 48, True, C_WHITE)
    add_textbox(
        s,
        Inches(0.7),
        Inches(2.6),
        Inches(11),
        Inches(0.5),
        "AI 여행 플래너 · Plan → Trip → After",
        22,
        False,
        RGBColor(0xBA, 0xE6, 0xFD),
    )
    add_textbox(
        s,
        Inches(0.7),
        Inches(3.4),
        Inches(11),
        Inches(0.8),
        "해외 MVP: 도쿄(기본) + 오사카(선택) · Expo / React Native · Gemini · Google Maps\n"
        "P0–P3 구현 완료 (P3: 9e57e3a) · 9ruDocs BlogDraft 호환",
        14,
        False,
        RGBColor(0xE0, 0xF2, 0xFE),
    )
    add_textbox(
        s,
        Inches(0.7),
        Inches(6.2),
        Inches(11),
        Inches(0.8),
        "제품 개요 문서  ·  github.com/sangyubkim/9ruTrip  ·  브랜드 #0c4a6e / #0369a1",
        12,
        False,
        RGBColor(0x7D, 0xD3, 0xFC),
    )

    # ========== 2. Vision ==========
    s = new_slide(prs)
    add_bg(s)
    add_header_bar(s, "비전 & 포지셔닝", "한 번의 여행 사이클을 앱에서 끝까지")
    add_card(s, Inches(0.45), Inches(1.25), Inches(12.4), Inches(1.5), C_SOFT, C_SOFT)
    add_textbox(
        s,
        Inches(0.7),
        Inches(1.4),
        Inches(12),
        Inches(1.2),
        "9ruTrip은 Gemini로 초기 일정을 만들고, 현장에서 동선·경비·기록을 관리한 뒤,\n"
        "여행 후 9ruDocs 호환 BlogDraft / WordPress로 발행까지 이어주는 모바일 여행 플래너입니다.",
        15,
        False,
        C_INK,
    )
    cols = [
        ("왜", ["해외 도시(도쿄/오사카) 빠른 MVP", "AI 일정 + 수동 DnD 하이브리드", "현장 가이드 알람·재루트", "경비 SMS + 계획 vs 실제"]),
        ("무엇", ["Expo Android-first 앱", "API(Gemini·Directions·WP)", "@9rutrip/shared 타입", "로컬 AsyncStorage 저장"]),
        ("누구", ["개인 여행자 / 블로거", "9ruDocs 발행 워크플로", "Android 실기기 우선", "국내(네이버)는 스캐폴드"]),
    ]
    for i, (h, items) in enumerate(cols):
        x = Inches(0.45 + i * 4.2)
        add_card(s, x, Inches(3.0), Inches(4.0), Inches(3.9))
        add_textbox(s, x + Inches(0.2), Inches(3.15), Inches(3.5), Inches(0.35), h, 16, True, C_PRIMARY)
        add_bullets(s, x + Inches(0.15), Inches(3.55), Inches(3.6), Inches(3.2), items, 12)

    # ========== 3. Plan / Trip / After ==========
    s = new_slide(prs)
    add_bg(s)
    add_header_bar(s, "여행 라이프사이클", "Plan → Trip(Active) → After")
    phases = [
        ("PLAN", C_PRIMARY, ["도시·박/일/인원 선택", "Gemini 일정 생성", "Day DnD · 카테고리 필터", "숙소 Top N · 교통 enrich", "수단 비교 · 지도+리스트"]),
        ("TRIP", C_ACCENT, ["status: active", "다음 액션 배너", "가이드 알람(로컬)", "AI 재루트 ON/OFF", "캡처 · 경비(SMS)"]),
        ("AFTER", RGBColor(0x07, 0x58, 0x85), ["계획 vs 실제 요약", "사진+리뷰 → Step", "export-draft (BlogDraft)", "WordPress 임시글/게시", "9ruDocs 파이프라인 연계"]),
    ]
    for i, (name, color, items) in enumerate(phases):
        x = Inches(0.45 + i * 4.2)
        head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), Inches(4.0), Inches(0.7))
        head.fill.solid()
        head.fill.fore_color.rgb = color
        head.line.fill.background()
        add_textbox(s, x, Inches(1.42), Inches(4.0), Inches(0.5), name, 20, True, C_WHITE, PP_ALIGN.CENTER)
        add_card(s, x, Inches(2.15), Inches(4.0), Inches(4.7))
        add_bullets(s, x + Inches(0.2), Inches(2.35), Inches(3.6), Inches(4.3), items, 13)
        if i < 2:
            add_textbox(s, x + Inches(3.85), Inches(1.4), Inches(0.5), Inches(0.5), "→", 22, True, C_MUTED, PP_ALIGN.CENTER)

    # ========== 4. P0–P3 ==========
    s = new_slide(prs)
    add_bg(s)
    add_header_bar(s, "기능 로드맵 P0–P3", "모두 구현됨 · 최신 P3 @ 9e57e3a")
    roadmap = [
        ("P0 MVP", [
            "Expo 셸 · 도쿄 일정",
            "Gemini + 폴백 코스",
            "Day 타임라인 DnD",
            "Maps 마커 · AsyncStorage",
            "캡처 · 현금 경비 · export-draft",
        ]),
        ("P1", [
            "가이드 알람 / 다음 액션",
            "AI 재루트 API",
            "SMS 경비 파싱",
            "WordPress 직접 발행",
            "교통 glance · 숙소 점수",
        ]),
        ("P2", [
            "Maps 키 배선 · Directions",
            "하버사인 폴백",
            "숙소 Top N + 재계산",
            "카테고리 DnD UX",
            "오사카 · Naver 스캐폴드",
        ]),
        ("P3", [
            "교통 수단 비교 UI/시트",
            "POST /trip/compare-transport",
            "transportOptions enrich",
            "Plan 지도(~37%)+리스트",
            "E2E 체크리스트 문서",
        ]),
    ]
    for i, (title, items) in enumerate(roadmap):
        x = Inches(0.35 + i * 3.2)
        add_card(s, x, Inches(1.25), Inches(3.05), Inches(5.7))
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.15), Inches(1.4), Inches(2.7), Inches(0.45))
        badge.fill.solid()
        badge.fill.fore_color.rgb = C_PRIMARY if i < 3 else C_ACCENT
        badge.line.fill.background()
        add_textbox(s, x + Inches(0.15), Inches(1.48), Inches(2.7), Inches(0.35), title, 13, True, C_WHITE, PP_ALIGN.CENTER)
        add_bullets(s, x + Inches(0.15), Inches(2.05), Inches(2.75), Inches(4.6), items, 11)

    # ========== 5. Design system ==========
    s = new_slide(prs)
    add_bg(s)
    add_header_bar(s, "UI 디자인 시스템", "앱 스크린 토큰 기준 · Malgun Gothic")
    # color swatches
    swatches = [
        ("#0c4a6e", C_PRIMARY, "Primary / Hero"),
        ("#0369a1", C_ACCENT, "Accent / CTA"),
        ("#0f172a", C_INK, "Ink / Title"),
        ("#f8fafc", C_BG, "Background"),
        ("#64748b", C_MUTED, "Muted / Meta"),
    ]
    for i, (hexv, rgb, label) in enumerate(swatches):
        x = Inches(0.5 + i * 2.5)
        sq = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.35), Inches(2.2), Inches(1.35))
        sq.fill.solid()
        sq.fill.fore_color.rgb = rgb
        sq.line.color.rgb = C_LINE
        tc = C_WHITE if i < 3 else C_INK
        add_textbox(s, x + Inches(0.1), Inches(1.55), Inches(2.0), Inches(0.4), hexv, 14, True, tc, PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.1), Inches(2.1), Inches(2.0), Inches(0.4), label, 11, False, tc, PP_ALIGN.CENTER)

    add_card(s, Inches(0.45), Inches(3.0), Inches(6.1), Inches(3.9))
    add_textbox(s, Inches(0.7), Inches(3.2), Inches(5.6), Inches(0.35), "타이포 · 컴포넌트", 16, True, C_PRIMARY)
    add_bullets(
        s,
        Inches(0.65),
        Inches(3.65),
        Inches(5.7),
        Inches(3.1),
        [
            "폰트: Malgun Gothic (본 문서) / RN 시스템 굵기 700–800",
            "Hero 카드: Primary 배경 + sky CTA (#38bdf8)",
            "탭/칩 ON: Accent 또는 Primary, OFF: slate-200",
            "리스트 행: 흰 카드 + e2e8f0 보더, 선택 시 sky-100",
            "완료 배지: emerald, 숙소 후보: orange 계열",
            "모달 시트: 흰 배경 + 핸들 바 + 모드 3열 비교",
        ],
        12,
    )
    add_card(s, Inches(6.8), Inches(3.0), Inches(6.05), Inches(3.9))
    add_textbox(s, Inches(7.05), Inches(3.2), Inches(5.6), Inches(0.35), "레이아웃 원칙", 16, True, C_PRIMARY)
    add_bullets(
        s,
        Inches(7.0),
        Inches(3.65),
        Inches(5.6),
        Inches(3.1),
        [
            "Home: 브랜드 히어로 → 여행 리스트",
            "Plan: Day 탭 + 필터 + 지도(~37%H) + DnD",
            "Active: 다음 액션 배너 상단 고정",
            "교통: glance 탭 → 비교 바텀시트",
            "After: 요약 스크롤 + 발행 CTA",
            "키 없음: Maps/Directions graceful degrade",
        ],
        12,
    )

    # ========== 6. Wireframes Home/Create/Plan ==========
    s = new_slide(prs)
    add_bg(s)
    add_header_bar(s, "와이어프레임 A", "Home · Create · Plan")
    wireframe_phone(
        s,
        Inches(0.4),
        Inches(1.2),
        "Home",
        [
            "##9ruTrip",
            "MVP · 도쿄 (Google Maps)",
            "[새 여행 만들기]",
            "API 설정",
            "---",
            "##저장된 여행",
            "도쿄 · 2박 3일",
            "2명 · planning · 계획 ¥…",
            "12곳 · 리뷰 3",
        ],
    )
    wireframe_phone(
        s,
        Inches(4.5),
        Inches(1.2),
        "Create",
        [
            "← 뒤로",
            "##여행 만들기",
            "해외 · Google Maps",
            "---",
            "도시  [도쿄] [오사카]",
            "박수 / 일수 / 인원",
            "[AI로 일정 생성]",
        ],
    )
    wireframe_phone(
        s,
        Inches(8.6),
        Inches(1.2),
        "Plan",
        [
            "Day1 Day2 Day3",
            "가이드알람 · AI재루트",
            "##지도 (~37%)",
            "마커 ↔ 리스트 연동",
            "---",
            "전체|맛집|관광|숙소",
            "≡ 장소 · glance 탭",
            "[지도][캡처][경비][요약]",
        ],
    )

    # ========== 7. Wireframes Transport/Active/After ==========
    s = new_slide(prs)
    add_bg(s)
    add_header_bar(s, "와이어프레임 B", "Transport 비교 · Active Trip · After")
    wireframe_phone(
        s,
        Inches(0.4),
        Inches(1.2),
        "Transport Compare",
        [
            "##이동 수단 비교",
            "→ 다음 장소명",
            "engine: Directions|haversine",
            "---",
            "[도보 ~12분 ¥0]",
            "[대중교통 ~8분 ¥…]",
            "[택시 ~5분 ¥…]",
            "선호 모드 저장",
            "glance / travelFromPrev*",
        ],
    )
    wireframe_phone(
        s,
        Inches(4.5),
        Inches(1.2),
        "Active Trip",
        [
            "##다음 액션 배너",
            "임박 Alert · 로컬 알림",
            "---",
            "Day 리스트 · 완료 체크",
            "이탈 시 AI 재루트",
            "[캡처] 사진+별점+문구",
            "[경비] 수동 / SMS 붙여넣기",
            "status: active",
        ],
    )
    wireframe_phone(
        s,
        Inches(8.6),
        Inches(1.2),
        "After / Summary",
        [
            "##비용 요약",
            "계획 vs 실제",
            "카테고리별 합계",
            "---",
            "[BlogDraft 내보내기]",
            "[WP 임시글]",
            "[WP 게시]",
            "리뷰 1개+ 필요",
            "9ruDocs Step 호환",
        ],
    )

    # ========== 8. Screen list ==========
    s = new_slide(prs)
    add_bg(s)
    add_header_bar(s, "화면 목록", "apps/mobile/src/screens · App.tsx 라우팅")
    screens = [
        ("HomeScreen", "여행 목록 · 생성/설정 진입"),
        ("CreateTripScreen", "도시(도쿄/오사카) · 박/일/인원 → itinerary"),
        ("PlanScreen", "Day DnD · 지도 · 비교시트 · 알람/재루트"),
        ("MapScreen", "전체 지도 뷰 (키 없으면 힌트)"),
        ("CaptureScreen", "카메라/갤러리 · PlaceReview"),
        ("ExpensesScreen", "수동 경비 · SMS 파싱 · share intent"),
        ("SummaryScreen", "비용 요약 · export-draft · WP publish"),
        ("SettingsScreen", "API Base URL · 헬스체크"),
    ]
    add_card(s, Inches(0.45), Inches(1.2), Inches(8.2), Inches(5.7))
    box = s.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(7.7), Inches(5.3))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (name, desc) in enumerate(screens):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        set_para(p, f"{name}", 13, True, C_PRIMARY, space_after=2)
        p2 = tf.add_paragraph()
        set_para(p2, f"    {desc}", 12, False, C_MUTED, space_after=8)

    add_card(s, Inches(8.9), Inches(1.2), Inches(3.95), Inches(5.7), C_SOFT, C_SOFT)
    add_textbox(s, Inches(9.1), Inches(1.4), Inches(3.5), Inches(0.4), "보조 컴포넌트", 14, True, C_PRIMARY)
    add_bullets(
        s,
        Inches(9.05),
        Inches(1.95),
        Inches(3.6),
        Inches(4.7),
        [
            "PlanDayMap",
            "TransportCompareSheet",
            "NextActionBanner",
            "useGuideAlarms",
            "trip API client",
            "cost / sms / nextAction utils",
        ],
        12,
    )

    # ========== 9. Architecture / API ==========
    s = new_slide(prs)
    add_bg(s)
    add_header_bar(s, "아키텍처 & API", "모노레포 · 포트 3011")
    # boxes
    nodes = [
        (0.5, "apps/mobile", "Expo SDK 55\nAndroid first\nAsyncStorage"),
        (4.7, "apps/api", "Gemini itinerary\nDirections/haversine\nWP publish"),
        (8.9, "packages/shared", "Trip / Step\nBlogDraft\n9ruDocs 호환"),
    ]
    for x, title, body in nodes:
        add_card(s, Inches(x), Inches(1.25), Inches(3.7), Inches(2.0))
        add_textbox(s, Inches(x + 0.15), Inches(1.35), Inches(3.4), Inches(0.35), title, 14, True, C_PRIMARY)
        add_textbox(s, Inches(x + 0.15), Inches(1.8), Inches(3.4), Inches(1.3), body, 12, False, C_MUTED)
    add_textbox(s, Inches(3.9), Inches(1.9), Inches(0.8), Inches(0.4), "→", 20, True, C_ACCENT, PP_ALIGN.CENTER)
    add_textbox(s, Inches(8.1), Inches(1.9), Inches(0.8), Inches(0.4), "→", 20, True, C_ACCENT, PP_ALIGN.CENTER)

    apis = [
        "GET  /health",
        "POST /trip/itinerary",
        "POST /trip/reroute",
        "POST /trip/enrich-transport",
        "POST /trip/compare-transport",
        "POST /trip/suggest-places",
        "POST /trip/parse-sms",
        "POST /trip/export-draft",
        "POST /wordpress/publish",
    ]
    add_card(s, Inches(0.45), Inches(3.5), Inches(12.4), Inches(3.4))
    add_textbox(s, Inches(0.7), Inches(3.65), Inches(11.8), Inches(0.35), "주요 엔드포인트", 14, True, C_PRIMARY)
    # three columns of APIs
    for col in range(3):
        chunk = apis[col * 3 : (col + 1) * 3]
        add_bullets(s, Inches(0.65 + col * 4.1), Inches(4.1), Inches(3.9), Inches(2.5), chunk, 12)

    # ========== 10. Gaps ==========
    s = new_slide(prs)
    add_bg(s)
    add_header_bar(s, "충분성 갭 (Sufficiency Gaps)", "MVP로 충분한 것 vs 아직 부족한 것")
    add_card(s, Inches(0.45), Inches(1.25), Inches(6.1), Inches(5.6), RGBColor(0xEC, 0xFD, 0xF5), RGBColor(0xA7, 0xF3, 0xD0))
    add_textbox(s, Inches(0.7), Inches(1.45), Inches(5.6), Inches(0.4), "충분한 영역 (P0–P3)", 16, True, C_OK)
    add_bullets(
        s,
        Inches(0.65),
        Inches(2.0),
        Inches(5.7),
        Inches(4.6),
        [
            "도쿄/오사카 AI 일정 + DnD 편집 루프",
            "교통 추정(Directions 또는 haversine) + 수단 비교",
            "현장 가이드 알람 · 재루트 · 캡처",
            "경비(수동/SMS)와 계획 vs 실제",
            "BlogDraft/WP 발행 경로",
            "키 없을 때 앱 크래시 없이 degrade",
            "실기기 E2E 체크리스트 문서화",
        ],
        12,
    )
    add_card(s, Inches(6.8), Inches(1.25), Inches(6.05), Inches(5.6), RGBColor(0xFF, 0xF7, 0xED), RGBColor(0xFE, 0xD7, 0xAA))
    add_textbox(s, Inches(7.05), Inches(1.45), Inches(5.6), Inches(0.4), "갭 / 리스크", 16, True, RGBColor(0x9A, 0x34, 0x12))
    add_bullets(
        s,
        Inches(7.0),
        Inches(2.0),
        Inches(5.6),
        Inches(4.6),
        [
            "국내 도시 Naver Maps 실연동 미완",
            "네이티브 SMS 인박스 자동 읽기 없음",
            "GPS 이탈(deviation) 힌트 미구현",
            "Routes API 고도화 · 실시간 혼잡 미반영",
            "iOS/웹은 2차 (Android first)",
            "공유 패키지 물리 병합 보류",
            "Maps/Gemini 키·비용 운영 의존",
        ],
        12,
    )

    # ========== 11. Next steps ==========
    s = new_slide(prs)
    add_bg(s)
    add_header_bar(s, "다음 단계", "README 「이후 (보류)」 기준")
    steps = [
        ("1", "국내 연동", "Naver Maps SDK로 국내 도시 실연동\n(현 mapProvider 스캐폴드 활용)"),
        ("2", "SMS 네이티브", "expo-dev-client + 권한 플러그인으로\n인박스 자동 파싱 경로"),
        ("3", "이탈 힌트", "GPS deviation 감지 → 재루트 제안\n(선택 기능)"),
        ("4", "플랫폼/품질", "Routes API 고도화 · iOS 패리티\n· 모노레포 물리 병합 검토"),
    ]
    for i, (num, title, body) in enumerate(steps):
        y = Inches(1.3 + i * 1.35)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55), y + Inches(0.15), Inches(0.7), Inches(0.7))
        circ.fill.solid()
        circ.fill.fore_color.rgb = C_PRIMARY
        circ.line.fill.background()
        add_textbox(s, Inches(0.55), y + Inches(0.28), Inches(0.7), Inches(0.45), num, 18, True, C_WHITE, PP_ALIGN.CENTER)
        add_card(s, Inches(1.5), y, Inches(11.3), Inches(1.2))
        add_textbox(s, Inches(1.75), y + Inches(0.15), Inches(10.8), Inches(0.35), title, 16, True, C_PRIMARY)
        add_textbox(s, Inches(1.75), y + Inches(0.5), Inches(10.8), Inches(0.55), body, 12, False, C_MUTED)

    # ========== 12. Closing ==========
    s = new_slide(prs)
    add_bg(s, C_PRIMARY)
    add_textbox(s, Inches(0.7), Inches(2.2), Inches(12), Inches(0.6), "요약", 28, True, C_WHITE)
    add_textbox(
        s,
        Inches(0.7),
        Inches(3.0),
        Inches(12),
        Inches(2.0),
        "9ruTrip P0–P3는 Plan·Trip·After 핵심 루프를 닫았습니다.\n"
        "다음 투자는 국내 맵·네이티브 SMS·이탈 힌트입니다.\n\n"
        "문서 재생성:  python docs/generate_overview_pptx.py",
        16,
        False,
        RGBColor(0xE0, 0xF2, 0xFE),
    )
    add_textbox(
        s,
        Inches(0.7),
        Inches(6.2),
        Inches(12),
        Inches(0.5),
        "https://github.com/sangyubkim/9ruTrip",
        13,
        False,
        RGBColor(0x7D, 0xD3, 0xFC),
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return len(prs.slides), OUT


if __name__ == "__main__":
    count, path = build()
    print(f"Wrote {path} ({count} slides, {path.stat().st_size} bytes)")
